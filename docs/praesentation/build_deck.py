#!/usr/bin/env python3
# Baut die Management-Praesentation (7 Folien, 16:9) fuer das
# VMware-Kapazitaets-Dashboard. Benoetigt python-pptx (Build-Werkzeug, NICHT
# Teil des Produkts). Screenshots werden aus SHOTS eingebettet.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import sys, os

SHOTS = sys.argv[1] if len(sys.argv) > 1 else "screenshots"
OUT   = sys.argv[2] if len(sys.argv) > 2 else "VMware-Kapazitaetsplanung.pptx"

# --- Farbpalette (an das Produkt angelehnt) ---
BG     = RGBColor(0x0B, 0x12, 0x20)   # dunkler Hintergrund
CARD   = RGBColor(0x15, 0x1E, 0x33)   # Panel
LINE   = RGBColor(0x2A, 0x36, 0x52)   # Rahmen
WHITE  = RGBColor(0xF2, 0xF5, 0xFA)
MUTED  = RGBColor(0x93, 0xA0, 0xB5)
ACCENT = RGBColor(0x7C, 0x83, 0xFF)   # Indigo (wie die Buttons)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
FONT   = "Segoe UI"

EMU_IN = 914400
TOTAL = 10           # Gesamtzahl der Folien (für die Fußzeile)
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: Liste von Absaetzen; jeder Absatz = Liste von (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT
    return tb


def picture(s, path, box_x, box_y, box_w, box_h, border=LINE):
    """Bild proportional in eine Box einpassen und zentrieren, mit Rahmen."""
    iw, ih = Image.open(path).size
    scale = min(box_w * EMU_IN / iw, box_h * EMU_IN / ih)
    w = int(iw * scale); h = int(ih * scale)
    x = int(box_x * EMU_IN + (box_w * EMU_IN - w) / 2)
    y = int(box_y * EMU_IN + (box_h * EMU_IN - h) / 2)
    pic = s.shapes.add_picture(path, Emu(x), Emu(y), width=Emu(w), height=Emu(h))
    pic.line.color.rgb = border; pic.line.width = Pt(1.25)
    return pic


def footer(s, n):
    rect(s, 0, 7.36, 13.333, 0.14, fill=ACCENT)
    text(s, 11.5, 7.02, 1.7, 0.3,
         [[("%d / %d" % (n, TOTAL), 10, MUTED, False)]], align=PP_ALIGN.RIGHT)
    text(s, 0.55, 7.02, 8, 0.3,
         [[("VMware Kapazitätsplanung", 10, MUTED, False)]])


def head(s, kicker, title):
    text(s, 0.6, 0.7, 12, 0.4, [[(kicker, 13, ACCENT, True)]])
    text(s, 0.6, 1.12, 12.1, 0.9, [[(title, 33, WHITE, True)]])


def bullets(items):
    out = []
    for it in items:
        out.append([("▸  ", 14, ACCENT, True), (it, 14, WHITE, False)])
    return out


def highlight(n, kicker, title, sub, items, img, page):
    s = slide()
    # linke Textspalte
    text(s, 0.6, 0.85, 5.2, 0.4, [[(kicker, 12.5, ACCENT, True)]])
    text(s, 0.6, 1.28, 4.9, 1.7, [[(title, 30, WHITE, True)]], line_spacing=1.02)
    text(s, 0.6, 2.9, 4.9, 0.9, [[(sub, 14.5, MUTED, False)]], line_spacing=1.1)
    text(s, 0.6, 3.95, 4.9, 3.0, bullets(items), space_after=11, line_spacing=1.05)
    # rechte Bildspalte auf einer Karte
    rect(s, 5.5, 0.9, 7.35, 5.75, fill=CARD, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    picture(s, os.path.join(SHOTS, img), 5.62, 1.02, 7.11, 5.51)
    footer(s, page)
    return s


# ---------------- Folie 1: Titel ----------------
s = slide()
rect(s, 0, 0, 0.28, 7.5, fill=ACCENT)
text(s, 0.9, 1.5, 11.5, 0.5, [[("KAPAZITÄTS-DASHBOARD FÜR VMWARE ARIA OPERATIONS",
                                14, ACCENT, True)]])
text(s, 0.86, 2.15, 11.6, 2.2,
     [[("Kapazität sehen.", 52, WHITE, True)],
      [("Planen. Freigeben.", 52, WHITE, True)]], line_spacing=1.0)
text(s, 0.9, 4.55, 10.8, 1.0,
     [[("Die gesamte virtuelle Landschaft auf einen Blick — in einem einzigen, "
        "abhängigkeitsfreien Werkzeug.", 19, MUTED, False)]], line_spacing=1.15)
rect(s, 0.92, 5.75, 6.4, 0.02, fill=LINE)
text(s, 0.9, 5.95, 11.8, 0.8,
     [[("Eine Python-Datei · keine Fremd-Abhängigkeiten · zweisprachig DE/EN · "
        "läuft auf jedem RHEL-Host · über 100 Releases", 13.5, MUTED, False)]],
     line_spacing=1.15)
footer(s, 1)

# ---------------- Folien 2–6: Top-5-Highlights ----------------
highlight(2, "HIGHLIGHT 1 · ÜBERBLICK",
          "Freie Kapazität auf einen Blick",
          "Sofort erkennen, wo Reserven sind — und wo es eng wird.",
          ["Freie vCPU, RAM & Storage je Cluster, farbcodiert nach Auslastung",
           "Ausfallreserve (N+1) und vSAN-Faktor bereits eingerechnet — keine Excel-Bastelei",
           "Ein Klick zum CSV-Export für Reports ans Management"],
          "01_kapazitaet.png", 2)

highlight(3, "HIGHLIGHT 2 · TIEFE",
          "Vom Cluster bis zur VM — ein Klick",
          "Volle Detailtiefe, ganz ohne vCenter-Login.",
          ["CPU, RAM, Storage, Netzwerk, Hosts & VMs je Cluster in einer Karte",
           "Tanzu-/Kubernetes-Reservierungen zählen automatisch mit",
           "Alle Zahlen mit derselben, nachvollziehbaren Rechenlogik"],
          "02_detail.png", 3)

highlight(4, "HIGHLIGHT 3 · GOVERNANCE",
          "Ressourcen mit klarem Freigabe-Prozess",
          "Nichts wird unbemerkt vergeben — alles ist nachvollziehbar.",
          ["Mehrstufiger Genehmigungs-Workflow je Team, inkl. Auto-Freigabe nach Schwellen",
           "Warnung, wenn eine Anfrage nicht mehr in den Cluster passt",
           "Lückenloses Audit-Log und automatische Mail-Benachrichtigungen"],
          "03_genehmigung.png", 4)

highlight(5, "HIGHLIGHT 4 · PROGNOSE",
          "Wohin wächst die Umgebung?",
          "Investitionsentscheidungen datenbasiert statt aus dem Bauch.",
          ["Trends aus bis zu 2 Jahren Historie: RAM/vCPU/Disk je VM, Auslastung, VM-Wachstum",
           "Frühzeitig erkennen, wann Erweiterungen fällig werden",
           "Selbst gezeichnete Charts — komplett offline, kein externer Dienst"],
          "05_statistik.png", 5)

highlight(6, "HIGHLIGHT 5 · REICHWEITE",
          "Alle Rechenzentren — auch die isolierten",
          "Eine Sicht für die ganze Landschaft, rollengerecht abgesichert.",
          ["Mehrere vROps-Quellen in einer Übersicht (Quellen-Badge je Cluster)",
           "Isolierte vCenter ohne Netz per Offline-Import (PowerCLI) einbinden",
           "Storage-Brücke: LUN-Erweiterungen direkt ans Storage-Team (API/CSV)"],
          "06_storage.png", 6)

# ---------------- Folie 7: Montage weiterer Funktionen ----------------
s = slide()
head(s, "MEHR FUNKTIONEN", "Und da steckt noch mehr drin")
montage = [
    ("07_vlan.png",       "VLAN-/Netzwerk-Suche",
     "Portgruppen & VLANs clusterübergreifend finden"),
    ("08_verwaltung.png", "Rollen, AD-Gruppen & Sichtbarkeit",
     "feingranular steuerbar, quellenbezogen filterbar"),
    ("09_apidocs.png",    "Eingebaute REST-API",
     "für Grafana, CMDB & Automatisierung"),
]
cx = 0.6
for img, cap_t, cap_s in montage:
    rect(s, cx, 1.75, 3.8, 3.45, fill=CARD, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    picture(s, os.path.join(SHOTS, img), cx + 0.12, 1.87, 3.56, 3.21)
    text(s, cx + 0.05, 5.34, 3.75, 0.4, [[(cap_t, 14, WHITE, True)]])
    text(s, cx + 0.05, 5.72, 3.75, 0.7, [[(cap_s, 11.5, MUTED, False)]],
         line_spacing=1.02)
    cx += 4.16
text(s, 0.6, 6.72, 12.1, 0.5,
     [[("Außerdem: ", 12.5, ACCENT, True),
       ("Hell-/Dunkel-Modus · SFTP-Backup · Mail-Benachrichtigungen & "
        "Erinnerungen · CSV-Import/Export · Auto-Freigabe · Reviewer-Handbuch · "
        "Ankündigungen · konfigurierbare Zeitzone · seitenweises Blättern",
        12.5, MUTED, False)]], line_spacing=1.05)
footer(s, 7)

# ---------------- Folie 8: Das Projekt in Zahlen ----------------
s = slide()
head(s, "ENTWICKLUNG", "Das Projekt in Zahlen")
text(s, 0.6, 1.9, 12, 0.5,
     [[("Von der ersten Zeile bis heute — in rund zwei Wochen aktiver "
        "Entwicklung.", 14.5, MUTED, False)]])
stats = [
    ("105", "Releases"), ("167", "Commits"),
    ("11.114", "Zeilen Code · 1 Datei"), ("0", "Fremd-Abhängigkeiten"),
    ("94", "automatisierte Smoke-Checks"), ("5", "Security-Scanner in der CI"),
    ("620+", "Übersetzungen · DE ↔ EN"), ("11", "Verwaltungs-Bereiche"),
]
tw, th, gx = 2.8, 1.9, 0.31
xs = [0.6 + i * (tw + gx) for i in range(4)]
for idx, (num, lbl) in enumerate(stats):
    x = xs[idx % 4]; y = 2.6 if idx < 4 else 4.65
    rect(s, x, y, tw, th, fill=CARD, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, y + 0.28, tw, 0.95, [[(num, 40, ACCENT, True)]],
         align=PP_ALIGN.CENTER)
    text(s, x + 0.15, y + 1.24, tw - 0.3, 0.6, [[(lbl, 12, MUTED, False)]],
         align=PP_ALIGN.CENTER, line_spacing=1.0)
footer(s, 8)

# ---------------- Folie 9: Fazit / Nutzenargumente ----------------
s = slide()
head(s, "FAZIT", "Warum das überzeugt")
cols = [
    ("Betrieb & Kosten", [
        "Eine einzige Python-Datei — kein Build, keine Datenbank, keine Fremd-Pakete",
        "Läuft auf jedem RHEL-Host; ein Update = eine Datei austauschen"]),
    ("Sicherheit & Compliance", [
        "AD-Anmeldung, rollenbasierte Sichtbarkeit, quellenbezogene Filter",
        "5 unabhängige Security-Scanner in der CI, geschützter main-Branch"]),
    ("Reichweite", [
        "Zweisprachig DE/EN, offline- und proxy-fähig",
        "Lesende API für Grafana, CMDB & Co."]),
    ("Pflege & Reife", [
        "Aktiv weiterentwickelt — über 100 Releases",
        "Reviewer-Handbuch, Architektur-Doku, Smoke-Tests inklusive"]),
]
positions = [(0.6, 2.35), (6.9, 2.35), (0.6, 4.55), (6.9, 4.55)]
for (title, items), (x, y) in zip(cols, positions):
    rect(s, x, y, 5.85, 1.95, fill=CARD, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + 0.32, y + 0.22, 5.3, 0.4, [[(title, 16, ACCENT, True)]])
    text(s, x + 0.32, y + 0.72, 5.25, 1.1, bullets(items),
         space_after=7, line_spacing=1.03)
text(s, 0.6, 6.72, 12, 0.5,
     [[("Wenig Angriffsfläche. Kein Vendor-Lock-in. Sofort einsatzbereit.",
        16, GREEN, True)]])
footer(s, 9)

# ---------------- Folie 10: Feature-Highlights (Sammlung) ----------------
s = slide()
head(s, "ALLES AUF EINEN BLICK", "Feature-Highlights")
groups = [
    ("Kapazität & Planung", [
        "Freie vCPU / RAM / Storage je Cluster",
        "N+1-Reserve & vSAN-Faktor eingerechnet",
        "Tanzu-/K8s-Namespaces zählen mit",
        "Workload aus vROps · 2-Jahres-Trends"]),
    ("Governance & Freigaben", [
        "Mehrstufiger Genehmigungs-Workflow",
        "Auto-Freigabe nach Schwellen",
        "Lückenloses Audit-Log",
        "Mail-Benachrichtigungen & Erinnerungen"]),
    ("Netzwerk & Storage", [
        "VLAN-/Portgruppen-Suche",
        "LUN-Drilldown inkl. NAA & WWPN",
        "Storage-Brücke ans Team (API/CSV)",
        "Mindest-LUN- & Namensfilter"]),
    ("Rollen & Sicherheit", [
        "AD-Anmeldung & AD-Gruppen",
        "Sichtbarkeits-Matrix je Rolle",
        "vROps-Quellen-Filter je Nutzer",
        "CSP-Härtung · 5 CI-Scanner"]),
    ("Integration & Betrieb", [
        "Multi-vROps + Offline-Import",
        "Lesende v1-API (Bearer-Token)",
        "SFTP-Backup · JSON oder SQLite",
        "/healthz · Proxy-fähig"]),
    ("Bedienung & Doku", [
        "Zweisprachig DE/EN · Hell-/Dunkel-Modus",
        "Sortieren, Blättern, Spalten wählen",
        "Konfigurierbare Zeitzone",
        "Reviewer-Handbuch & Architektur-Doku"]),
]
bw, bh, gxx = 3.84, 2.45, 0.3
bxs = [0.6, 0.6 + bw + gxx, 0.6 + 2 * (bw + gxx)]
for idx, (title, items) in enumerate(groups):
    x = bxs[idx % 3]; y = 1.95 if idx < 3 else 4.55
    rect(s, x, y, bw, bh, fill=CARD, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + 0.26, y + 0.18, bw - 0.4, 0.4, [[(title, 13.5, ACCENT, True)]])
    runs = [[("•  ", 10.5, GREEN, True), (it, 10.5, WHITE, False)] for it in items]
    text(s, x + 0.26, y + 0.64, bw - 0.46, bh - 0.75, runs,
         space_after=4.5, line_spacing=1.0)
footer(s, 10)

prs.save(OUT)
print("gespeichert:", OUT, "-", len(prs.slides._sldIdLst), "Folien")
