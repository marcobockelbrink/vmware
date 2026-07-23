#!/usr/bin/env python3
# Kleine Präsentation zu den lustigen Codenamen + Sticker-Logos.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import sys, os

PNG = sys.argv[1] if len(sys.argv) > 1 else "logos"
OUT = sys.argv[2] if len(sys.argv) > 2 else "Codename-Vorschlaege.pptx"

BG    = RGBColor(0x0B, 0x12, 0x20)
CARD  = RGBColor(0x15, 0x1E, 0x33)
LINE  = RGBColor(0x2A, 0x36, 0x52)
WHITE = RGBColor(0xF2, 0xF5, 0xFA)
MUTED = RGBColor(0x93, 0xA0, 0xB5)
ACC   = RGBColor(0x7C, 0x83, 0xFF)
GREEN = RGBColor(0x34, 0xD3, 0x99)
FONT  = "Segoe UI"
EMU = 914400
TOTAL = 9

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s


def rect(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         sa=6, ls=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0); p.line_spacing = ls
        for (t, sz, c, b) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = FONT
    return tb


def logo(s, path, box_x, box_y, box_w, box_h):
    iw, ih = Image.open(path).size
    sc = min(box_w * EMU / iw, box_h * EMU / ih)
    w = int(iw * sc); h = int(ih * sc)
    x = int(box_x * EMU + (box_w * EMU - w) / 2)
    y = int(box_y * EMU + (box_h * EMU - h) / 2)
    s.shapes.add_picture(path, Emu(x), Emu(y), width=Emu(w), height=Emu(h))


def footer(s, n):
    rect(s, 0, 7.36, 13.333, 0.14, fill=ACC)
    text(s, 11.5, 7.02, 1.7, 0.3, [[("%d / %d" % (n, TOTAL), 10, MUTED, False)]],
         align=PP_ALIGN.RIGHT)
    text(s, 0.55, 7.02, 8, 0.3, [[("Codename-Vorschläge · VMware Kapazitätsplanung", 10, MUTED, False)]])


# ---- Folie 1: Titel ----
s = slide()
rect(s, 0, 0, 0.28, 7.5, fill=ACC)
text(s, 0.9, 1.7, 11.5, 0.5, [[("PROJEKT-CODENAME GESUCHT", 14, ACC, True)]])
text(s, 0.86, 2.25, 11.6, 1.6, [[("Die lustige Runde 🎉", 50, WHITE, True)]])
text(s, 0.9, 3.9, 11.0, 1.0,
     [[("7 Namen · 7 Sticker-Logos · frisch aus der Design-Werkstatt.",
        20, MUTED, False)]])
text(s, 0.9, 5.7, 11.8, 0.6,
     [[("Alle mit Virtualisierungs- & Kapa-Bezug — und bereit für die "
        "Aufkleber-Maschine.", 14, MUTED, False)]])
footer(s, 1)

# ---- Folien 2–8: je ein Name ----
names = [
 ("1_vacancy.png", "vAcancy",
  "„Vacancy“ (freie Zimmer) im VMware-v-Look.",
  "Zeigt buchstäblich, wo noch Platz frei ist — und klingt fast wie ein echtes VMware-Produkt."),
 ("2_platzhirsch.png", "Platzhirsch",
  "„Platz“ = Kapazität — der Platzhirsch sichert sich das beste Revier.",
  "Wer bekommt die vCPUs? Der Platzhirsch entscheidet. Der Charme-Sieger fürs Team."),
 ("3_clustertetris.png", "ClusterTetris",
  "Kapazitätsplanung ist am Ende Tetris.",
  "Passt der Block (VM) noch in den Cluster — oder ist oben Schluss?"),
 ("4_rambo.png", "RAMbo",
  "Der Muskel hinter deinem Speicher.",
  "„First Blood“ gibt's, wenn der RAM alle ist."),
 ("5_overbooked.png", "Overbooked",
  "Wie die Airline beim Overbooking — VMware nennt es „Overcommit“.",
  "Warnt dich, bevor die Maschine (VM) am Boden bleibt."),
 ("6_clusterphobie.png", "Clusterphobie",
  "Die Angst, dass der Cluster volläuft.",
  "Und diese App ist die Therapie dagegen."),
 ("7_luftballons.png", "99 Luftballons",
  "„Ballooning“ ist ein echter VMware-Speichertrick.",
  "Die Nena-Referenz gibt's gratis obendrauf."),
]
for i, (img, name, gag, pitch) in enumerate(names):
    s = slide()
    logo(s, os.path.join(PNG, img), 0.7, 1.35, 4.9, 4.9)
    text(s, 6.2, 1.7, 6.5, 0.4, [[("VORSCHLAG %d / 7" % (i + 1), 13, ACC, True)]])
    text(s, 6.2, 2.15, 6.6, 1.1, [[(name, 46, WHITE, True)]])
    text(s, 6.2, 3.5, 6.5, 0.4, [[("Der Gag", 13, GREEN, True)]])
    text(s, 6.2, 3.9, 6.5, 1.1, [[(gag, 17, WHITE, False)]], ls=1.1)
    text(s, 6.2, 5.05, 6.5, 0.4, [[("Warum es zündet", 13, GREEN, True)]])
    text(s, 6.2, 5.45, 6.5, 1.1, [[(pitch, 15, MUTED, False)]], ls=1.12)
    footer(s, i + 2)

# ---- Folie 9: Sticker-Sheet ----
s = slide()
text(s, 0.6, 0.6, 12, 0.4, [[("DAS STICKER-SHEET", 13, ACC, True)]])
text(s, 0.6, 1.0, 12, 0.8, [[("Bereit für die Aufkleber-Maschine 🎉", 32, WHITE, True)]])
allpng = [n[0] for n in names]
sz = 2.55
row1 = allpng[:4]; row2 = allpng[4:]
xs1 = [0.6 + i * (sz + 0.56) for i in range(4)]
for x, img in zip(xs1, row1):
    logo(s, os.path.join(PNG, img), x, 2.0, sz, sz)
start2 = (13.333 - (len(row2) * sz + (len(row2) - 1) * 0.56)) / 2
xs2 = [start2 + i * (sz + 0.56) for i in range(len(row2))]
for x, img in zip(xs2, row2):
    logo(s, os.path.join(PNG, img), x, 4.55, sz, sz)
text(s, 0.6, 6.95, 12.1, 0.4,
     [[("Welcher wird's? ", 14, GREEN, True),
       ("Dein Call — dann webe ich ihn ins Dashboard, die README und die Präsentation ein.",
        14, MUTED, False)]])
footer(s, 9)

prs.save(OUT)
print("gespeichert:", OUT, "-", len(prs.slides._sldIdLst), "Folien")
